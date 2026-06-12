from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta Advanx ──────────────────────────────────────────────────────────────
NEGRO      = RGBColor(0x0D, 0x0D, 0x0D)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
ACENTO     = RGBColor(0x00, 0xB4, 0xD8)   # azul eléctrico
ACENTO2    = RGBColor(0x02, 0x3E, 0x8A)   # azul oscuro
GRIS       = RGBColor(0x2B, 0x2D, 0x42)   # gris oscuro
GRIS_CLARO = RGBColor(0xEF, 0xEF, 0xEF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completamente en blanco

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=24, bold=False, color=BLANCO,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  size=18, color=BLANCO, bold_first=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

def bg(slide, color=NEGRO):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=color)

def accent_bar(slide, top=0.72):
    add_rect(slide, 0, top, 13.33, 0.06, fill_color=ACENTO)

def slide_title(slide, title, subtitle=None):
    add_text(slide, title, 0.6, 0.18, 12, 0.55,
             size=28, bold=True, color=ACENTO)
    accent_bar(slide, top=0.75)
    if subtitle:
        add_text(slide, subtitle, 0.6, 0.82, 12, 0.45,
                 size=16, color=GRIS_CLARO)

def fase_badge(slide, label, l, t, color=ACENTO2):
    add_rect(slide, l, t, 2.8, 0.38, fill_color=color)
    add_text(slide, label, l+0.08, t+0.03, 2.64, 0.32,
             size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NEGRO)
add_rect(s, 0, 0, 0.18, 7.5, fill_color=ACENTO)          # barra lateral
add_rect(s, 0.18, 3.2, 13.15, 0.06, fill_color=ACENTO2)  # línea central

add_text(s, "ADVANX", 0.5, 1.2, 12, 1.4,
         size=96, bold=True, color=BLANCO)
add_text(s, "Ordena.  Mide.  Escala.", 0.5, 2.5, 12, 0.6,
         size=28, color=ACENTO)
add_text(s,
    '"Las PyMEs en México no mueren por falta de clientes.\nMueren porque no pueden con los que ya tienen."',
    0.5, 3.4, 11, 1.2, size=20, italic=True, color=GRIS_CLARO)
add_text(s, "Sesión de alineación · Junio 2026", 0.5, 6.8, 12, 0.45,
         size=14, color=ACENTO, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — El problema
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "El problema que resolvemos")

puntos = [
    "No sabe qué hay en inventario sin preguntarle a alguien",
    "Un reporte le toma 2 horas hacerlo a mano",
    "Cuando un empleado clave se va, el conocimiento se va con él",
    "Las decisiones se toman por intuición, no por datos",
]
for i, p in enumerate(puntos):
    add_rect(s, 0.6, 1.1 + i*1.3, 0.06, 0.9, fill_color=ACENTO)
    add_text(s, p, 0.85, 1.12 + i*1.3, 11.5, 0.9,
             size=20, color=BLANCO)

add_rect(s, 0.6, 6.2, 12.1, 0.7, fill_color=ACENTO2)
add_text(s,
    "Síntoma visible: alguien tiene como trabajo principal actualizar un Excel todos los días.",
    0.75, 6.25, 11.8, 0.6, size=18, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Cliente ideal
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "¿A quién le vendemos?")

criterios = [
    ("Tamaño",              "20–80 empleados"),
    ("Industrias",          "Manufactura · Distribución · Alimentos · Servicios con operación física"),
    ("Síntoma visible",     "Alguien actualiza Excel a diario como su trabajo principal"),
    ("Dolor del dueño",     '"No sé qué pasa en mi negocio sin preguntarle a alguien"'),
    ("Filtro clave",        "¿No puedo con los clientes que tengo? → SÍ es nuestro cliente"),
]
for i, (lbl, val) in enumerate(criterios):
    y = 1.1 + i * 1.1
    add_rect(s, 0.6, y, 3.2, 0.75, fill_color=ACENTO2)
    add_text(s, lbl, 0.68, y+0.1, 3.0, 0.55, size=15, bold=True, color=BLANCO)
    add_text(s, val, 4.0, y+0.08, 9.0, 0.65, size=16, color=GRIS_CLARO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Qué es Advanx
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "¿Qué es Advanx?")

add_text(s,
    "Diagnostica, ordena y automatiza los procesos clave de una PyME\npara que pueda crecer con datos reales.",
    0.6, 1.0, 12, 0.9, size=22, color=BLANCO)

add_rect(s, 0.6, 2.1, 12.1, 1.1, fill_color=ACENTO)
add_text(s, "El diferenciador:", 0.85, 2.18, 4, 0.4, size=16, bold=True, color=NEGRO)
add_text(s, "No hacemos reportes. Construimos sistemas.", 0.85, 2.55, 11.5, 0.55,
         size=26, bold=True, color=NEGRO)

versus = [
    ("Contratar TI interno",  "Advanx trae método + ejecución, no solo el recurso"),
    ("Implementar un ERP",    "Primero ordenamos el proceso, luego lo digitalizamos"),
    ("Seguir con Excel",      "El costo de la ineficiencia ya superó el costo del cambio"),
    ("Consultora grande",     "Acceso directo al experto · Precio de PyME, no de corporativo"),
]
add_text(s, "Versus las alternativas", 0.6, 3.4, 6, 0.4, size=15, bold=True, color=ACENTO)
for i, (alt, razon) in enumerate(versus):
    y = 3.85 + i * 0.78
    add_rect(s, 0.6, y, 3.5, 0.6, fill_color=GRIS)
    add_text(s, alt,   0.68, y+0.08, 3.3, 0.45, size=14, bold=True, color=ACENTO)
    add_text(s, razon, 4.25, y+0.1,  8.7, 0.5,  size=14, color=GRIS_CLARO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Vista general 4 fases
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "La metodología: 4 fases · 16 sesiones · ~4 meses")

fases = [
    ("FASE 1\nPlaneación",    "¿Qué eres y\nhacia dónde vas?",  "Base\nEstratégica",    ACENTO2),
    ("FASE 2\nOrganización",  "¿Cómo operas y\nquién hace qué?","Manual de\nOperaciones",GRIS),
    ("FASE 3\nDirección",     "¿Cómo mides\ny decides?",        "Kit de\nDecisiones",   ACENTO2),
    ("FASE 4\nControl",       "¿Cómo creces\nsolo?",            "Playbook\nCompleto",   GRIS),
]

box_w = 2.9
gap   = 0.28
start = 0.5

for i, (titulo, pregunta, doc, color) in enumerate(fases):
    x = start + i * (box_w + gap)
    # caja principal
    add_rect(s, x, 1.05, box_w, 3.6, fill_color=color)
    add_text(s, titulo,   x+0.12, 1.15, box_w-0.2, 0.85,
             size=16, bold=True, color=ACENTO if color == GRIS else BLANCO,
             align=PP_ALIGN.CENTER)
    add_rect(s, x, 1.98, box_w, 0.04, fill_color=ACENTO)
    add_text(s, "4 sesiones", x+0.12, 2.08, box_w-0.2, 0.4,
             size=13, color=GRIS_CLARO, align=PP_ALIGN.CENTER)
    add_text(s, pregunta, x+0.12, 2.5, box_w-0.2, 1.0,
             size=14, color=BLANCO, align=PP_ALIGN.CENTER, italic=True)
    # flecha entre cajas
    if i < 3:
        add_text(s, "→", x + box_w + 0.02, 2.6, gap + 0.05, 0.5,
                 size=22, bold=True, color=ACENTO, align=PP_ALIGN.CENTER)
    # documento integrador
    add_rect(s, x, 4.85, box_w, 0.65, fill_color=ACENTO)
    add_text(s, doc, x+0.08, 4.9, box_w-0.12, 0.55,
             size=14, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)

add_text(s, "Documento integrador de cada fase  ↑",
         0.5, 5.6, 12.3, 0.4, size=13, color=GRIS_CLARO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Fase 1 detalle
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Fase 1 — Planeación", "¿Qué eres, para quién, cuánto cuesta y hacia dónde vas?")

sesiones = [
    ("S1", "Propuesta de valor + ICP",         "Mensaje claro que no cambia según con quién hablas"),
    ("S2", "Motivador de compra real",           "Por qué te compran de verdad — no lo que asumes"),
    ("S3", "Estructura de costos + obj. SMART", "Punto de equilibrio en unidades · metas con fecha"),
    ("S4", "Criterios de crecimiento + BMC",    "Cuándo y cómo escalar · síntesis visual del negocio"),
]
for i, (num, tema, resultado) in enumerate(sesiones):
    y = 1.15 + i * 1.28
    add_rect(s, 0.5, y, 0.7, 0.7, fill_color=ACENTO)
    add_text(s, num, 0.5, y+0.1, 0.7, 0.5, size=18, bold=True,
             color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(s, tema,      1.35, y,      5.5, 0.42, size=17, bold=True, color=BLANCO)
    add_text(s, resultado, 1.35, y+0.42, 5.5, 0.45, size=15, color=GRIS_CLARO)
    add_rect(s, 7.1, y, 0.04, 0.72, fill_color=ACENTO2)
    add_text(s, "→ Entregable", 7.25, y+0.05, 1.5, 0.35, size=12,
             color=ACENTO, italic=True)

add_rect(s, 0.5, 6.25, 12.3, 0.75, fill_color=ACENTO2)
add_text(s, "Documento integrador: BASE ESTRATÉGICA — propuesta de valor + ICP + costos + objetivos SMART",
         0.65, 6.32, 12.0, 0.6, size=16, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Fase 2 detalle
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Fase 2 — Organización", "¿Cómo hacemos las cosas y quién es responsable de qué?")

sesiones = [
    ("S5", "Mapeo de procesos clave",       "Mapa de cómo fluye el trabajo de punta a punta"),
    ("S6", "Documentación de SOPs",          "Pasos escritos para los procesos que más afectan al cliente"),
    ("S7", "Roles y accountability",         "Quién es responsable de qué — sin depender de la memoria de nadie"),
    ("S8", "Revisión de implementación",     "Ajuste tras 2 semanas ejecutando en la realidad"),
]
for i, (num, tema, resultado) in enumerate(sesiones):
    y = 1.15 + i * 1.28
    add_rect(s, 0.5, y, 0.7, 0.7, fill_color=ACENTO)
    add_text(s, num, 0.5, y+0.1, 0.7, 0.5, size=18, bold=True,
             color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(s, tema,      1.35, y,      9.5, 0.42, size=17, bold=True, color=BLANCO)
    add_text(s, resultado, 1.35, y+0.42, 9.5, 0.45, size=15, color=GRIS_CLARO)

add_rect(s, 0.5, 6.25, 12.3, 0.75, fill_color=ACENTO2)
add_text(s, "Documento integrador: MANUAL DE OPERACIONES — mapa de procesos + SOPs + roles",
         0.65, 6.32, 12.0, 0.6, size=16, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Fase 3 detalle
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Fase 3 — Dirección", "¿Cómo sé si voy bien y qué hago cuando algo no va?")

sesiones = [
    ("S9",  "Definición de KPIs",            "3–5 métricas que realmente importan (no vanity metrics)"),
    ("S10", "Dashboard simple",               "Google Sheet que se actualiza en 10 min/semana"),
    ("S11", "Protocolo de decisiones",        "Reglas escritas: si X baja de Y → hago Z"),
    ("S12", "Primera revisión con datos",     "Primera decisión tomada con datos reales — en vivo"),
]
for i, (num, tema, resultado) in enumerate(sesiones):
    y = 1.15 + i * 1.28
    add_rect(s, 0.5, y, 0.7, 0.7, fill_color=ACENTO)
    add_text(s, num, 0.5, y+0.1, 0.7, 0.5, size=18, bold=True,
             color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(s, tema,      1.35, y,      9.5, 0.42, size=17, bold=True, color=BLANCO)
    add_text(s, resultado, 1.35, y+0.42, 9.5, 0.45, size=15, color=GRIS_CLARO)

add_rect(s, 0.5, 6.25, 12.3, 0.75, fill_color=ACENTO2)
add_text(s, "Documento integrador: KIT DE DECISIONES — dashboard + KPIs + protocolo de decisiones",
         0.65, 6.32, 12.0, 0.6, size=16, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Fase 4 detalle
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Fase 4 — Control", "¿Cómo hago que el crecimiento se sostenga solo?")

sesiones = [
    ("S13", "Diseño del flywheel",            "Mapa del ciclo de crecimiento propio del negocio"),
    ("S14", "Sistema de retención y referidos","Mecanismo concreto que activa el flywheel"),
    ("S15", "Criterios de escalamiento",       "Cuándo contratar, abrir turno, subir precios — con datos"),
    ("S16", "Documento maestro + roadmap",     "Todo consolidado + hoja de ruta para el siguiente año"),
]
for i, (num, tema, resultado) in enumerate(sesiones):
    y = 1.15 + i * 1.28
    add_rect(s, 0.5, y, 0.7, 0.7, fill_color=ACENTO)
    add_text(s, num, 0.5, y+0.1, 0.7, 0.5, size=18, bold=True,
             color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(s, tema,      1.35, y,      9.5, 0.42, size=17, bold=True, color=BLANCO)
    add_text(s, resultado, 1.35, y+0.42, 9.5, 0.45, size=15, color=GRIS_CLARO)

add_rect(s, 0.5, 6.25, 12.3, 0.75, fill_color=ACENTO2)
add_text(s, "Documento integrador: PLAYBOOK COMPLETO — consolida las 3 fases + flywheel + roadmap 12 meses",
         0.65, 6.32, 12.0, 0.6, size=16, bold=True, color=BLANCO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Caso UMA Spa
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Caso real: UMA Spa", "Primer engagement de Advanx · Fase 1 en curso")

# Columna izquierda — Antes
add_rect(s, 0.5, 1.05, 5.8, 4.8, fill_color=GRIS)
add_text(s, "ANTES", 0.65, 1.12, 5.5, 0.45, size=16, bold=True, color=ACENTO)
antes = [
    "Sin propuesta de valor consistente",
    "Sin saber por qué le compraban",
    "Precios definidos por intuición",
    "Sin estructura de costos",
    "Sin punto de equilibrio conocido",
]
for i, linea in enumerate(antes):
    add_text(s, "✗  " + linea, 0.7, 1.62 + i*0.72, 5.4, 0.6,
             size=15, color=GRIS_CLARO)

# Columna derecha — Hallazgos
add_rect(s, 6.6, 1.05, 6.2, 4.8, fill_color=ACENTO2)
add_text(s, "HALLAZGOS EN 3 SESIONES", 6.75, 1.12, 6.0, 0.45,
         size=16, bold=True, color=BLANCO)

hallazgos = [
    "Los clientes no compran un masaje — compran\nla certeza de presencia total sin prisa.",
    "Punto de equilibrio: 34 masajes/mes",
    "Capacidad instalada: 96 masajes/mes",
    "Mes típico: 27 masajes",
    "Brecha a rentabilidad: 7 masajes más\n(< 2 por semana)",
]
for i, h in enumerate(hallazgos):
    add_text(s, "→  " + h, 6.75, 1.62 + i*0.84, 5.9, 0.75,
             size=14, color=BLANCO)

add_rect(s, 0.5, 6.05, 12.3, 0.85, fill_color=ACENTO)
add_text(s,
    "El problema no es capacidad. Es visibilidad. Ese es el tipo de hallazgo que solo da la metodología.",
    0.65, 6.12, 12.0, 0.7, size=17, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Modelo de negocio
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Modelo de negocio")

servicios = [
    ("Diagnóstico",           "$25,000 MXN",  "2–3 semanas", "Entrada de bajo compromiso. El cliente ve el mapa antes de comprometerse."),
    ("Sprint 1 mes",          "$22,000 MXN",  "4 semanas",   "1 proceso automatizado de punta a punta. Para quien quiere ver resultados rápido."),
    ("Acompañamiento\n3 meses","$18,000/mes", "12 semanas",  "El engagement completo. 1 cliente activo cubre la meta mínima de ingreso."),
]
for i, (nombre, precio, dur, desc) in enumerate(servicios):
    x = 0.5 + i * 4.28
    add_rect(s, x, 1.05, 3.98, 4.5, fill_color=GRIS)
    add_rect(s, x, 1.05, 3.98, 0.05, fill_color=ACENTO)
    add_text(s, nombre, x+0.15, 1.15, 3.7, 0.7,
             size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(s, precio, x+0.1, 1.95, 3.8, 0.65,
             size=28, bold=True, color=ACENTO, align=PP_ALIGN.CENTER)
    add_text(s, dur,    x+0.1, 2.65, 3.8, 0.45,
             size=15, color=GRIS_CLARO, align=PP_ALIGN.CENTER)
    add_text(s, desc,   x+0.15, 3.18, 3.7, 1.3,
             size=14, color=GRIS_CLARO, align=PP_ALIGN.CENTER)

add_rect(s, 0.5, 5.75, 12.3, 0.55, fill_color=ACENTO2)
add_text(s,
    "Gancho: si el cliente contrata acompañamiento justo después del diagnóstico → $10,000 de crédito.",
    0.65, 5.8, 12.0, 0.45, size=15, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# Math
add_text(s, "1 cliente en acompañamiento = $18,000/mes = meta mínima cubierta",
         0.5, 6.45, 12.3, 0.5, size=16, bold=True, color=ACENTO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Mercado
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Por qué el momento es ahora")

datos = [
    ("7.7%",  "crecimiento anual\nde consultoría PyME"),
    ("12.8%", "crecimiento anual\nde consultoría tecnológica"),
    ("76%",   "de PyMEs mexicanas\naún sin digitalizar"),
]
for i, (num, label) in enumerate(datos):
    x = 0.7 + i * 4.1
    add_rect(s, x, 1.1, 3.6, 2.4, fill_color=ACENTO2)
    add_text(s, num,   x+0.1, 1.25, 3.4, 1.1,
             size=52, bold=True, color=ACENTO, align=PP_ALIGN.CENTER)
    add_text(s, label, x+0.1, 2.35, 3.4, 0.9,
             size=15, color=BLANCO, align=PP_ALIGN.CENTER)

add_rect(s, 0.5, 3.75, 12.3, 1.5, fill_color=GRIS)
add_text(s, "El espacio integrado en el Valle de Toluca está vacío.",
         0.65, 3.82, 12.0, 0.55, size=22, bold=True, color=ACENTO)
add_text(s,
    "Nadie combina las 4 dimensiones: procesos + personas + datos + IA con presencia local.\n"
    "Los competidores atacan cada eje por separado. El espacio integrado no tiene dueño.",
    0.65, 4.4, 12.0, 0.8, size=16, color=GRIS_CLARO)

add_rect(s, 0.5, 5.45, 12.3, 1.5, fill_color=ACENTO)
add_text(s,
    "La demanda existe. El mercado es inmaduro.\nQuien llega primero con metodología probada define el estándar.",
    0.65, 5.55, 12.0, 1.2, size=20, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Por qué nosotros
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Por qué nosotros")

razones = [
    ("Caso real validado",      "UMA Spa — metodología ejecutada en campo, no diseñada en papel"),
    ("Sin intermediarios",      "El cliente habla directo con quien implementa y entrega"),
    ("Enfoque correcto",        "Primero ordenamos, luego automatizamos — no se automatiza el caos"),
    ("Precio justo",            "Precio de PyME, velocidad de startup, calidad de consultoría senior"),
]
for i, (titulo, desc) in enumerate(razones):
    y = 1.15 + i * 1.35
    add_rect(s, 0.5, y, 0.08, 0.9, fill_color=ACENTO)
    add_text(s, titulo, 0.75, y,      11.5, 0.45, size=19, bold=True, color=ACENTO)
    add_text(s, desc,   0.75, y+0.45, 11.5, 0.55, size=16, color=GRIS_CLARO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conversación con la socia
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "La conversación que necesitamos tener")

preguntas = [
    "¿Qué capacidades traes tú que complementan las mías?",
    "¿Cómo dividimos roles dentro de un engagement con un cliente?",
    "¿Cuál sería el primer cliente que abordaríamos juntos?",
    "¿Qué necesitas ver o entender antes de comprometerte?",
]
for i, p in enumerate(preguntas):
    y = 1.2 + i * 1.35
    add_rect(s, 0.5, y, 0.9, 0.9, fill_color=ACENTO)
    add_text(s, str(i+1), 0.5, y+0.12, 0.9, 0.65,
             size=30, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)
    add_text(s, p, 1.55, y+0.18, 11.0, 0.65, size=20, color=BLANCO)

add_rect(s, 0.5, 6.6, 12.3, 0.65, fill_color=ACENTO2)
add_text(s, "Advanx · Ordena. Mide. Escala.",
         0.65, 6.68, 12.0, 0.5, size=16, bold=True, color=ACENTO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Guardar
# ══════════════════════════════════════════════════════════════════════════════
output = r"C:\Users\1544\Documents\GOOGLE DRIVE RECOVERY\OBSIDIAN\DIGITAL BRAIN\Advanx — Presentación para Socia.pptx"
prs.save(output)
print(f"✓ Guardado en: {output}")
